"""
Tools for the AI Agent (Sandboxed Version)
Each tool is protected by safety guardrails and runs in a per-session sandbox.
"""
import os
import subprocess
import sys
import json
import time
import re
import shutil
import httpx
import trafilatura

try:
    from duckduckgo_search import DDGS
    HAS_DDG = True
except ImportError:
    HAS_DDG = False

from config import (
    WORKSPACE_DIR,
    MAX_FILE_READ_BYTES,
    MAX_TOOL_OUTPUT_CHARS,
)
from safety import guard, SafetyViolation
from sandbox_session import session_manager


def get_sandbox_status(session_id: str = None) -> dict:
    """Public helper for the UI: describe the current sandbox mode."""
    if session_id:
        try:
            sb = session_manager.get_or_create(session_id)
            return sb.get_status()
        except Exception:
            pass
    return {"mode": "unknown", "available": False}


# ============================================================
# TOOL IMPLEMENTATIONS
# ============================================================

def _truncate(text: str, max_chars: int = MAX_TOOL_OUTPUT_CHARS) -> str:
    """Truncate long outputs and append a notice."""
    if not text:
        return text or ""
    if len(text) <= max_chars:
        return text
    return (
        text[:max_chars]
        + f"\n\n... [truncated {len(text) - max_chars} chars; "
          "use a more specific tool (grep, head, tail, edit_file) to see more]"
    )


def run_bash(command: str, timeout: int = 30, session_id: str = None) -> dict:
    """Run a shell command in the session's sandbox."""
    check = guard.validate_command(command)
    if not check["safe"]:
        return {
            "status": "blocked",
            "output": check["message"],
            "risk_level": check["risk_level"],
        }

    timeout = min(max(timeout, 1), 120)
    sandbox = session_manager.get_or_create(session_id)
    result = sandbox.run_command(command, timeout)

    # Add safety metadata
    result["risk_level"] = check["risk_level"]
    if check["risk_level"] == "risky":
        result["output"] = check["message"] + "\n\n---\n\n" + result.get("output", "")

    result["output"] = _truncate(result.get("output", ""))
    return result


def read_file(path: str, max_bytes: int = MAX_FILE_READ_BYTES, offset: int = 0,
              session_id: str = None) -> dict:
    """Read a file with safety checks. Supports offset for large files."""
    check = guard.validate_file_read(path)
    if not check["safe"]:
        return {"status": "blocked", "output": check["message"], "risk_level": "blocked"}

    try:
        resolved = guard.validate_path(path, allow_workspace_only=False, must_exist=True)

        if os.path.isdir(resolved):
            return {
                "status": "error",
                "output": f"'{resolved}' is a directory, not a file. Use list_files instead.",
            }

        # Binary-file guard
        binary_extensions = {
            '.docx', '.xlsx', '.pptx', '.doc', '.xls', '.ppt',
            '.pdf', '.zip', '.tar', '.gz', '.rar', '.7z',
            '.exe', '.dll', '.so', '.dylib',
            '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.ico',
            '.mp3', '.mp4', '.avi', '.mov', '.wav', '.flac',
            '.pyc', '.pyo', '.class', '.o', '.obj',
        }
        _, ext = os.path.splitext(resolved.lower())
        if ext in binary_extensions:
            file_size = os.path.getsize(resolved)
            return {
                "status": "error",
                "output": (
                    f"Cannot read binary file: {resolved}\n"
                    f"File type: {ext}\n"
                    f"File size: {file_size} bytes\n\n"
                    f"💡 Use `view_file` tool to preview this file, "
                    f"or use `run_python` with appropriate libraries."
                ),
            }

        file_size = os.path.getsize(resolved)
        if file_size > max_bytes + offset:
            with open(resolved, 'rb') as f:
                f.seek(offset)
                raw = f.read(max_bytes)
            content = raw.decode('utf-8', errors='replace')
            content = (
                f"[Showing bytes {offset}..{offset+max_bytes} of {file_size}. "
                f"Pass offset= to read more.]\n\n" + content
            )
        else:
            with open(resolved, 'r', encoding='utf-8', errors='replace') as f:
                if offset:
                    f.seek(offset)
                content = f.read()

        # Heuristic binary sniff
        if content.count('\x00') > 10:
            return {
                "status": "error",
                "output": (
                    f"File appears to be binary: {resolved}\n"
                    f"File size: {file_size} bytes\n"
                    f"Null bytes detected. Use `view_file` or `run_python`."
                ),
            }

        return {
            "status": "success",
            "output": _truncate(content),
            "file_size": file_size,
            "truncated": file_size > max_bytes + offset,
        }
    except SafetyViolation as e:
        return {"status": "blocked", "output": str(e), "risk_level": "blocked"}
    except UnicodeDecodeError:
        file_size = os.path.getsize(resolved) if os.path.exists(resolved) else 0
        return {
            "status": "error",
            "output": (
                f"Cannot decode file as text: {resolved}\n"
                f"File size: {file_size} bytes\n\n"
                f"💡 This file appears to be binary. Use `view_file` or `run_python`."
            ),
        }
    except Exception as e:
        return {"status": "error", "output": str(e)}


def view_file(path: str, session_id: str = None) -> dict:
    """
    View/preview a file. Auto-detects file type and returns a preview.
    Supports: docx, pdf, pptx, xlsx, images, text files.
    """
    check = guard.validate_file_read(path)
    if not check["safe"]:
        return {"status": "blocked", "output": check["message"], "risk_level": "blocked"}

    try:
        resolved = guard.validate_path(path, allow_workspace_only=False, must_exist=True)

        if os.path.isdir(resolved):
            return {
                "status": "error",
                "output": f"'{resolved}' is a directory. Use `list_files` instead.",
            }

        _, ext = os.path.splitext(resolved.lower())
        file_size = os.path.getsize(resolved)

        # Route to appropriate handler based on extension
        if ext in ('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'):
            return _view_image(resolved, ext, file_size)

        elif ext == '.pdf':
            return _view_pdf(resolved, file_size, session_id)

        elif ext == '.docx':
            return _view_docx(resolved, file_size, session_id)

        elif ext == '.pptx':
            return _view_pptx(resolved, file_size, session_id)

        elif ext in ('.xlsx', '.xls'):
            return _view_excel(resolved, file_size, session_id)

        elif ext in ('.csv',):
            return _view_csv(resolved, file_size)

        else:
            # Fall back to text read
            return read_file(path, session_id=session_id)

    except SafetyViolation as e:
        return {"status": "blocked", "output": str(e), "risk_level": "blocked"}
    except Exception as e:
        return {"status": "error", "output": f"view_file error: {type(e).__name__}: {e}"}


def _view_image(path: str, ext: str, file_size: int) -> dict:
    """Return image metadata (actual preview is handled by the UI)."""
    try:
        from PIL import Image
        img = Image.open(path)
        width, height = img.size
        return {
            "status": "success",
            "output": (
                f"🖼️ Image: {os.path.basename(path)}\n"
                f"Size: {width}x{height} pixels\n"
                f"Mode: {img.mode}\n"
                f"File size: {file_size:,} bytes\n"
                f"Format: {ext[1:].upper()}\n\n"
                f"📁 Path: {path}"
            ),
            "file_type": "image",
            "path": path,
            "width": width,
            "height": height,
        }
    except ImportError:
        return {
            "status": "success",
            "output": (
                f"🖼️ Image: {os.path.basename(path)}\n"
                f"File size: {file_size:,} bytes\n"
                f"Format: {ext[1:].upper()}\n\n"
                f"📁 Path: {path}"
            ),
            "file_type": "image",
            "path": path,
        }
    except Exception as e:
        return {"status": "error", "output": f"Cannot read image: {e}"}


def _view_pdf(path: str, file_size: int, session_id: str = None) -> dict:
    """Extract text from PDF using the session sandbox."""
    sandbox = session_manager.get_or_create(session_id)
    safe_path = repr(path)
    code = f"""
import sys
try:
    import pdfplumber
    with pdfplumber.open({safe_path}) as pdf:
        pages = pdf.pages[:20]
        for i, page in enumerate(pages):
            text = page.extract_text() or ''
            print(f'--- Page {{i+1}} ---')
            print(text)
            print()
except ImportError:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader({safe_path})
        for i, page in enumerate(reader.pages[:20]):
            text = page.extract_text() or ''
            print(f'--- Page {{i+1}} ---')
            print(text)
            print()
    except Exception as e:
        print(f'Error: {{e}}')
"""
    result = sandbox.run_python(code)
    result["file_type"] = "pdf"
    result["path"] = path
    result["file_size"] = file_size
    result["output"] = _truncate(result.get("output", ""), max_chars=15000)
    return result


def _view_docx(path: str, file_size: int, session_id: str = None) -> dict:
    """Extract text from DOCX using the session sandbox."""
    sandbox = session_manager.get_or_create(session_id)
    safe_path = repr(path)
    code = f"""
try:
    from docx import Document
    doc = Document({safe_path})
    for para in doc.paragraphs:
        print(para.text)
    for i, table in enumerate(doc.tables):
        print(f'\\n--- Table {{i+1}} ---')
        for row in table.rows:
            print(' | '.join(cell.text for cell in row.cells))
except Exception as e:
    print(f'Error reading docx: {{e}}')
"""
    result = sandbox.run_python(code)
    result["file_type"] = "docx"
    result["path"] = path
    result["file_size"] = file_size
    result["output"] = _truncate(result.get("output", ""), max_chars=15000)
    return result


def _view_pptx(path: str, file_size: int, session_id: str = None) -> dict:
    """Extract text from PPTX using the session sandbox."""
    sandbox = session_manager.get_or_create(session_id)
    safe_path = repr(path)
    code = f"""
try:
    from pptx import Presentation
    prs = Presentation({safe_path})
    for i, slide in enumerate(prs.slides):
        print(f'--- Slide {{i+1}} ---')
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                print(shape.text)
        print()
except Exception as e:
    print(f'Error reading pptx: {{e}}')
"""
    result = sandbox.run_python(code)
    result["file_type"] = "pptx"
    result["path"] = path
    result["file_size"] = file_size
    result["output"] = _truncate(result.get("output", ""), max_chars=15000)
    return result


def _view_excel(path: str, file_size: int, session_id: str = None) -> dict:
    """Extract data from Excel using the session sandbox."""
    sandbox = session_manager.get_or_create(session_id)
    safe_path = repr(path)
    code = f"""
try:
    import openpyxl
    wb = openpyxl.load_workbook({safe_path}, read_only=True, data_only=True)
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        print(f'--- Sheet: {{sheet_name}} ---')
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= 50:
                print(f'... ({{ws.max_row - 50}} more rows)')
                break
            print(' | '.join(str(cell) if cell is not None else '' for cell in row))
        print()
    wb.close()
except Exception as e:
    print(f'Error reading excel: {{e}}')
"""
    result = sandbox.run_python(code)
    result["file_type"] = "excel"
    result["path"] = path
    result["file_size"] = file_size
    result["output"] = _truncate(result.get("output", ""), max_chars=15000)
    return result


def _view_csv(path: str, file_size: int) -> dict:
    """Read CSV file (text-based, no sandbox needed)."""
    try:
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            lines = []
            for i, line in enumerate(f):
                if i >= 100:
                    lines.append(f"... (more rows)")
                    break
                lines.append(line.rstrip())
        return {
            "status": "success",
            "output": "\n".join(lines),
            "file_type": "csv",
            "path": path,
            "file_size": file_size,
        }
    except Exception as e:
        return {"status": "error", "output": f"Cannot read CSV: {e}"}


def write_file(path: str, content: str, session_id: str = None) -> dict:
    """Write a file with safety checks (workspace only)."""
    check = guard.validate_file_write(path)
    if not check["safe"]:
        return {"status": "blocked", "output": check["message"], "risk_level": "blocked"}

    try:
        resolved = guard.validate_path(path, allow_workspace_only=True)

        parent_dir = os.path.dirname(resolved)
        os.makedirs(parent_dir, exist_ok=True)

        with open(resolved, 'w', encoding='utf-8') as f:
            f.write(content)

        try:
            display_path = os.path.relpath(resolved, WORKSPACE_DIR)
        except ValueError:
            display_path = resolved

        return {
            "status": "success",
            "output": f"✅ File written: {display_path} ({len(content)} bytes)\nLocation: {resolved}",
        }
    except SafetyViolation as e:
        return {"status": "blocked", "output": str(e), "risk_level": "blocked"}
    except Exception as e:
        return {"status": "error", "output": str(e)}


def edit_file(path: str, old_string: str, new_string: str, replace_all: bool = False,
              session_id: str = None) -> dict:
    """
    Surgical edit: replace `old_string` with `new_string` in `path`.
    Fails if old_string is not found, or (when replace_all=False) appears more than once.
    Workspace only.
    """
    check = guard.validate_file_write(path)
    if not check["safe"]:
        return {"status": "blocked", "output": check["message"], "risk_level": "blocked"}

    if not old_string:
        return {"status": "error", "output": "old_string cannot be empty"}

    try:
        resolved = guard.validate_path(path, allow_workspace_only=True, must_exist=True)

        with open(resolved, 'r', encoding='utf-8') as f:
            content = f.read()

        occurrences = content.count(old_string)
        if occurrences == 0:
            return {
                "status": "error",
                "output": (
                    f"old_string not found in {resolved}.\n"
                    f"💡 Read the file first to see the exact text (whitespace matters)."
                ),
            }
        if occurrences > 1 and not replace_all:
            return {
                "status": "error",
                "output": (
                    f"old_string appears {occurrences} times in {resolved}.\n"
                    f"💡 Provide more surrounding context to make it unique, "
                    f"or set replace_all=true to replace every occurrence."
                ),
            }

        new_content = (
            content.replace(old_string, new_string)
            if replace_all
            else content.replace(old_string, new_string, 1)
        )

        with open(resolved, 'w', encoding='utf-8') as f:
            f.write(new_content)

        try:
            display_path = os.path.relpath(resolved, WORKSPACE_DIR)
        except ValueError:
            display_path = resolved

        return {
            "status": "success",
            "output": (
                f"✅ Edited {display_path}\n"
                f"Replaced {occurrences if replace_all else 1} occurrence(s). "
                f"({len(old_string)} -> {len(new_string)} chars)"
            ),
        }
    except SafetyViolation as e:
        return {"status": "blocked", "output": str(e), "risk_level": "blocked"}
    except Exception as e:
        return {"status": "error", "output": str(e)}


def list_files(path: str = ".", session_id: str = None) -> dict:
    """List files in workspace."""
    try:
        resolved = guard.validate_path(path, allow_workspace_only=True, must_exist=True)
        items = []
        for item in sorted(os.listdir(resolved)):
            full_path = os.path.join(resolved, item)
            item_type = "📁 DIR " if os.path.isdir(full_path) else "📄 FILE"
            size = os.path.getsize(full_path) if os.path.isfile(full_path) else "-"
            items.append(f"  {item_type} {item} ({size})")
        header = f"Contents of: {os.path.relpath(resolved, WORKSPACE_DIR) or 'workspace'}\n"
        return {
            "status": "success",
            "output": header + ("\n".join(items) if items else "(empty directory)"),
        }
    except SafetyViolation as e:
        return {"status": "blocked", "output": str(e), "risk_level": "blocked"}
    except Exception as e:
        return {"status": "error", "output": str(e)}


def web_search(query: str, max_results: int = 5, session_id: str = None) -> dict:
    """Search the web. Uses duckduckgo_search (real results) if installed, else falls back."""
    if not HAS_DDG:
        return {
            "status": "error",
            "output": (
                "duckduckgo-search is not installed. Run:\n"
                "  pip install duckduckgo-search==6.3.5\n"
                "and restart the agent."
            ),
        }
    try:
        results = []
        with DDGS() as ddgs:
            for r in ddgs.text(query, max_results=max_results):
                title = r.get("title", "")
                href = r.get("href", "")
                body = r.get("body", "")
                results.append(f"### {title}\n{body}\nURL: {href}\n")
        if not results:
            results.append("No results. Try a different query.")
        return {"status": "success", "output": _truncate("\n".join(results), max_chars=6000)}
    except Exception as e:
        return {"status": "error", "output": f"Search error: {e}"}


def fetch_page(url: str, max_chars: int = 12000, session_id: str = None) -> dict:
    """Fetch a URL and return the main text content (trafilatura-extracted)."""
    if not (url.startswith("http://") or url.startswith("https://")):
        return {"status": "error", "output": "URL must start with http:// or https://"}
    try:
        with httpx.Client(
            timeout=20,
            follow_redirects=True,
            headers={"User-Agent": "Mozilla/5.0 (compatible; ai-agent/1.0)"},
        ) as client:
            resp = client.get(url)
            resp.raise_for_status()
            html = resp.text

        extracted = trafilatura.extract(html, include_links=False, include_images=False) or ""
        if not extracted.strip():
            # Fallback: very crude text-only extraction
            extracted = re.sub(r"<script.*?</script>", "", html, flags=re.S | re.I)
            extracted = re.sub(r"<style.*?</style>", "", extracted, flags=re.S | re.I)
            extracted = re.sub(r"<[^>]+>", " ", extracted)
            extracted = re.sub(r"\s+", " ", extracted).strip()

        head = f"Source: {url}\nStatus: {resp.status_code}\nLength: {len(extracted)} chars\n\n"
        return {
            "status": "success",
            "output": _truncate(head + extracted, max_chars=max_chars),
        }
    except httpx.HTTPError as e:
        return {"status": "error", "output": f"HTTP error: {e}"}
    except Exception as e:
        return {"status": "error", "output": f"Fetch error: {e}"}


def pip_install(package: str, session_id: str = None) -> dict:
    """
    Install a Python package in the session's sandbox (temporary).
    The package will be available for the rest of the session only.
    It will NOT be installed on the user's host machine.
    """
    check = guard.validate_pip_install(package)
    if not check["safe"]:
        return {"status": "blocked", "output": check["message"], "risk_level": "blocked"}

    sandbox = session_manager.get_or_create(session_id)
    result = sandbox.install_package(package)

    if result.get("status") == "success":
        result["output"] = (
            f"✅ Installed '{package}' in sandbox (temporary)\n"
            f"📦 Available for the rest of this session only.\n"
            f"🧹 Will be cleaned up when session ends.\n"
            f"🛡️ Host machine is NOT modified."
        )
    else:
        result["output"] = f"❌ Failed to install '{package}'\n" + result.get("output", "")

    return result


def run_python(code: str, session_id: str = None) -> dict:
    """Execute Python code in the session's sandbox."""
    check = guard.validate_python_code(code)
    if not check["safe"]:
        return {"status": "blocked", "output": check["message"], "risk_level": "blocked"}

    sandbox = session_manager.get_or_create(session_id)
    result = sandbox.run_python(code)

    # Add safety metadata
    if check["risk_level"] == "risky":
        result["output"] = (
            check["message"]
            + "\n\n---\n\n"
            + "⚠️ This is a warning, not a block. The code ran.\n\n"
            + result.get("output", "")
        )

    result["risk_level"] = check["risk_level"]
    result["output"] = _truncate(result.get("output", ""))
    return result


# ============================================================
# IMAGE / FILE TOOLS
# ============================================================

def download_file(url: str, filename: str = None, max_bytes: int = 50_000_000,
                  session_id: str = None) -> dict:
    """
    Download a file from a URL into the workspace.
    Used to grab logos, images, PDFs, etc. from the web.
    """
    if not (url.startswith("http://") or url.startswith("https://")):
        return {"status": "error", "output": "URL must start with http:// or https://"}
    try:
        # Derive a safe filename
        if not filename:
            from urllib.parse import urlparse
            path = urlparse(url).path
            filename = os.path.basename(path) or "downloaded_file"
            # Strip query/fragment residue
            filename = filename.split("?")[0].split("#")[0]
        # Sanitize: only allow basename, no path traversal
        filename = os.path.basename(filename)
        if not filename:
            filename = "downloaded_file"
        dest = guard.validate_path(filename, allow_workspace_only=True)

        with httpx.Client(
            timeout=60,
            follow_redirects=True,
            headers={"User-Agent": "Mozilla/5.0 (compatible; ai-agent/1.0)"},
        ) as client:
            with client.stream("GET", url) as resp:
                resp.raise_for_status()
                # Content-length check (best effort)
                cl = resp.headers.get("content-length")
                if cl and int(cl) > max_bytes:
                    return {"status": "error", "output": f"File too large: {cl} bytes (max {max_bytes})"}
                with open(dest, "wb") as f:
                    downloaded = 0
                    for chunk in resp.iter_bytes(chunk_size=65536):
                        downloaded += len(chunk)
                        if downloaded > max_bytes:
                            f.close()
                            os.remove(dest)
                            return {"status": "error", "output": f"File too large (> {max_bytes} bytes)"}
                        f.write(chunk)

        size = os.path.getsize(dest)
        try:
            display = os.path.relpath(dest, WORKSPACE_DIR)
        except ValueError:
            display = dest
        return {
            "status": "success",
            "output": f"✅ Downloaded: {display} ({size:,} bytes)\nFrom: {url}",
            "path": display,
            "size": size,
        }
    except SafetyViolation as e:
        return {"status": "blocked", "output": str(e), "risk_level": "blocked"}
    except httpx.HTTPError as e:
        return {"status": "error", "output": f"HTTP error: {e}"}
    except Exception as e:
        return {"status": "error", "output": f"Download error: {e}"}


def image_search(query: str, max_results: int = 5, session_id: str = None) -> dict:
    """
    Search the web for images. Uses duckduckgo_search if available.
    Returns URLs of images matching the query. Then use download_file to grab them.
    """
    if not HAS_DDG:
        return {
            "status": "error",
            "output": "duckduckgo-search is not installed. Run: pip install duckduckgo-search",
        }
    try:
        results = []
        with DDGS() as ddgs:
            for r in ddgs.images(query, max_results=max_results):
                title = r.get("title", "")
                image_url = r.get("image") or r.get("url", "")
                thumbnail = r.get("thumbnail", "")
                width = r.get("width", "")
                height = r.get("height", "")
                results.append(
                    f"### {title}\n"
                    f"Image URL: {image_url}\n"
                    f"Thumbnail: {thumbnail}\n"
                    f"Size: {width}x{height}\n"
                )
        if not results:
            results.append("No images found. Try a different query.")
        return {"status": "success", "output": _truncate("\n".join(results), max_chars=6000)}
    except Exception as e:
        return {"status": "error", "output": f"Image search error: {e}"}


def process_image(
    path: str,
    output: str = None,
    resize: tuple = None,
    crop: tuple = None,
    convert_to: str = None,
    make_transparent: bool = False,
    session_id: str = None,
) -> dict:
    """
    Process an image using Pillow: resize, crop, convert format,
    or attempt to make a white background transparent.
    """
    try:
        from PIL import Image

        resolved = guard.validate_path(path, allow_workspace_only=True, must_exist=True)
        if not output:
            output = os.path.basename(resolved)
        out_path = guard.validate_path(output, allow_workspace_only=True)

        img = Image.open(resolved)

        original_size = img.size
        original_mode = img.mode

        if crop:
            img = img.crop(crop)
        if resize:
            img = img.resize(resize, Image.LANCZOS)
        if make_transparent:
            img = img.convert("RGBA")
            pixels = img.load()
            threshold = 240
            for y in range(img.height):
                for x in range(img.width):
                    r, g, b, a = pixels[x, y]
                    if r >= threshold and g >= threshold and b >= threshold:
                        pixels[x, y] = (r, g, b, 0)
        if convert_to:
            target_format = convert_to.upper()
            if target_format in ("JPEG", "JPG") and img.mode in ("RGBA", "LA", "P"):
                bg = Image.new("RGB", img.size, (255, 255, 255))
                bg.paste(img, mask=img.split()[-1])
                img = bg
            img.save(out_path, format=target_format)
        else:
            img.save(out_path)

        new_size = os.path.getsize(out_path)
        try:
            display = os.path.relpath(out_path, WORKSPACE_DIR)
        except ValueError:
            display = out_path
        ops = []
        if crop: ops.append(f"crop{crop}")
        if resize: ops.append(f"resize{resize}")
        if make_transparent: ops.append("transparent-bg")
        if convert_to: ops.append(f"->{convert_to}")
        op_str = ", ".join(ops) or "saved"
        return {
            "status": "success",
            "output": (
                f"✅ Processed: {display}\n"
                f"Operations: {op_str}\n"
                f"Original: {original_size[0]}x{original_size[1]} {original_mode}\n"
                f"New size: {new_size:,} bytes"
            ),
            "path": display,
        }
    except SafetyViolation as e:
        return {"status": "blocked", "output": str(e), "risk_level": "blocked"}
    except ImportError:
        return {"status": "error", "output": "Pillow not installed in sandbox. Try: pip_install('Pillow')"}
    except Exception as e:
        return {"status": "error", "output": f"Image error: {type(e).__name__}: {e}"}


def remove_background(path: str, output: str = None, session_id: str = None) -> dict:
    """
    Remove the background from an image, leaving the foreground subject
    with a transparent background. Uses the `rembg` library if installed,
    otherwise falls back to a simple white-to-transparent conversion.
    """
    try:
        resolved = guard.validate_path(path, allow_workspace_only=True, must_exist=True)
        if not output:
            base, _ = os.path.splitext(os.path.basename(resolved))
            output = f"{base}_nobg.png"
        out_path = guard.validate_path(output, allow_workspace_only=True)

        # Try rembg first (uses AI model, much better quality)
        try:
            from rembg import remove as rembg_remove
            from PIL import Image
            with open(resolved, "rb") as f:
                input_data = f.read()
            output_data = rembg_remove(input_data)
            with open(out_path, "wb") as f:
                f.write(output_data)
            try:
                display = os.path.relpath(out_path, WORKSPACE_DIR)
            except ValueError:
                display = out_path
            return {
                "status": "success",
                "output": f"✅ Removed background (using rembg AI): {display} ({os.path.getsize(out_path):,} bytes)",
                "path": display,
            }
        except ImportError:
            # Fall back to threshold-based transparency
            return process_image(
                path=path,
                output=output,
                make_transparent=True,
                convert_to="PNG",
                session_id=session_id,
            )

    except SafetyViolation as e:
        return {"status": "blocked", "output": str(e), "risk_level": "blocked"}
    except Exception as e:
        return {"status": "error", "output": f"BG-remove error: {type(e).__name__}: {e}"}


# ============================================================
# TOOL DEFINITIONS (OpenAI function-calling format)
# ============================================================

TOOL_DEFINITIONS = [
    {
        "type": "function",
        "function": {
            "name": "run_bash",
            "description": (
                "Run a shell command in the sandboxed workspace directory. "
                "DANGEROUS COMMANDS (format, del on system files, shutdown, "
                "rm -rf /, etc.) will be BLOCKED. Use this for git, npm, "
                "compiling, running programs, etc."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "command": {"type": "string", "description": "The shell command to execute."},
                    "timeout": {"type": "integer", "description": "Timeout in seconds (max 120)."},
                },
                "required": ["command"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_file",
            "description": (
                "Read a file. Supports an offset/max_bytes for large files. "
                "Binary files and credential files are blocked. "
                "Use view_file for binary files (docx, pdf, images, etc.)."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "File path."},
                    "max_bytes": {"type": "integer", "description": "Max bytes to read (default 50000)."},
                    "offset": {"type": "integer", "description": "Byte offset to start from."},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "view_file",
            "description": (
                "View/preview a file. Auto-detects file type and returns a preview. "
                "Supports: docx, pdf, pptx, xlsx, images (png/jpg/gif/bmp/webp), "
                "csv, and text files. Use this instead of read_file for binary files."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "File path to view."},
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "write_file",
            "description": (
                "Create or overwrite a file. ONLY works within the workspace directory. "
                "Use edit_file for surgical changes to existing files."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "File path (within workspace)."},
                    "content": {"type": "string", "description": "The full content to write."},
                },
                "required": ["path", "content"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "edit_file",
            "description": (
                "Surgically edit a file by replacing a specific string. "
                "Fails if the string is not found, or appears more than once "
                "(unless replace_all=true). This is much cheaper than write_file "
                "for small changes. Workspace only."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "File path (within workspace)."},
                    "old_string": {"type": "string", "description": "The exact string to replace."},
                    "new_string": {"type": "string", "description": "What to replace it with."},
                    "replace_all": {
                        "type": "boolean",
                        "description": "Replace every occurrence (default false).",
                        "default": False,
                    },
                },
                "required": ["path", "old_string", "new_string"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "list_files",
            "description": "List files in a directory within the workspace.",
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Directory path (default: workspace root)."},
                },
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "web_search",
            "description": (
                "Search the web using DuckDuckGo. Returns real results "
                "(title, snippet, URL). Use this for current information, "
                "documentation, or facts."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "The search query."},
                    "max_results": {"type": "integer", "description": "How many results (default 5)."},
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "fetch_page",
            "description": (
                "Fetch a URL and return the main text content (article body). "
                "Use after web_search to read a specific page."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "The URL to fetch (http/https)."},
                    "max_chars": {"type": "integer", "description": "Max chars to return (default 12000)."},
                },
                "required": ["url"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "pip_install",
            "description": (
                "Install a Python package in the session's TEMPORARY sandbox. "
                "The package will be available for the rest of the session only "
                "and will NOT be installed on the user's host machine. "
                "Known malicious/typosquat packages are blocked."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "package": {
                        "type": "string",
                        "description": "Package name (e.g. 'requests', 'numpy==1.26.0').",
                    },
                },
                "required": ["package"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "run_python",
            "description": (
                "Execute Python code in the session's sandbox. "
                "Dangerous patterns (eval, exec, os.system) are FLAGGED but "
                "the code still runs - this is a warning, not a sandbox. "
                "Use for calculations, data processing, quick scripts."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "code": {"type": "string", "description": "The Python code to execute."},
                },
                "required": ["code"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "download_file",
            "description": (
                "Download a file from a URL into the workspace. "
                "Useful for grabbing logos, images, PDFs, etc. from the web. "
                "After downloading, you can process it with process_image or "
                "remove_background."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "url": {"type": "string", "description": "HTTP/HTTPS URL to download."},
                    "filename": {
                        "type": "string",
                        "description": "Optional output filename (defaults to URL basename).",
                    },
                },
                "required": ["url"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "image_search",
            "description": (
                "Search the web for images matching a query. Returns URLs. "
                "Use download_file afterward to save a chosen image to the workspace. "
                "Best for finding official logos, photos, icons."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "query": {"type": "string", "description": "What to search for."},
                    "max_results": {"type": "integer", "description": "How many images (default 5)."},
                },
                "required": ["query"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "process_image",
            "description": (
                "Process an image: resize, crop, convert format, or make a "
                "near-white background transparent. Uses Pillow. "
                "Common usage: resize a downloaded logo to fit in a header."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Input image path (workspace)."},
                    "output": {"type": "string", "description": "Output filename (default: overwrite)."},
                    "resize": {
                        "type": "array",
                        "items": {"type": "integer"},
                        "description": "Optional (width, height) to resize to.",
                    },
                    "crop": {
                        "type": "array",
                        "items": {"type": "integer"},
                        "description": "Optional (left, top, right, bottom) crop box.",
                    },
                    "convert_to": {
                        "type": "string",
                        "description": "Optional output format: PNG, JPEG, WEBP, etc.",
                    },
                    "make_transparent": {
                        "type": "boolean",
                        "description": "If true, convert near-white pixels to transparent.",
                        "default": False,
                    },
                },
                "required": ["path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "remove_background",
            "description": (
                "Remove the background from an image, leaving the foreground "
                "subject with a transparent background. Uses rembg (AI model) "
                "if installed; falls back to white-to-transparent conversion. "
                "Saves as PNG."
            ),
            "parameters": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Input image path (workspace)."},
                    "output": {"type": "string", "description": "Output filename (default: <name>_nobg.png)."},
                },
                "required": ["path"],
            },
        },
    },
]


TOOL_FUNCTIONS = {
    "run_bash": run_bash,
    "read_file": read_file,
    "view_file": view_file,
    "write_file": write_file,
    "edit_file": edit_file,
    "list_files": list_files,
    "web_search": web_search,
    "fetch_page": fetch_page,
    "pip_install": pip_install,
    "run_python": run_python,
    "download_file": download_file,
    "image_search": image_search,
    "process_image": process_image,
    "remove_background": remove_background,
}
